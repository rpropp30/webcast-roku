# -*- coding: utf-8 -*-
"""
Generator for: Powerline Worker Electrical Theory & Energized-Work Safety Training
Builds a polished 16:9 PowerPoint deck (~78 slides) with a consistent theme.

Run:  python3 build_presentation.py
Out:  Powerline_Electrical_Theory_and_Safety_Training.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os
import re
import json
from PIL import Image

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
with open(os.path.join(os.path.dirname(os.path.abspath(__file__)), "notes.json"), encoding="utf-8") as _nf:
    NOTES = json.load(_nf)


def _norm(s):
    s = re.sub(r"<[^>]+>", "", s)
    s = s.replace("&amp;", "&").replace("&rsquo;", "'").replace("&ldquo;", "").replace("&rdquo;", "")
    s = s.replace("“", "").replace("”", "").replace("‘", "'").replace("’", "'")
    s = re.sub(r"\s+", " ", s).strip()
    return s

# ----------------------------------------------------------------------------
# Palette
# ----------------------------------------------------------------------------
NAVY      = RGBColor(0x0B, 0x25, 0x45)   # deep navy (primary dark)
NAVY2     = RGBColor(0x13, 0x40, 0x74)   # mid blue
STEEL     = RGBColor(0x3A, 0x5A, 0x80)   # steel blue
AMBER     = RGBColor(0xF5, 0xB1, 0x00)   # safety amber (accent)
AMBER_D   = RGBColor(0xC8, 0x8A, 0x00)   # darker amber
DANGER    = RGBColor(0xD4, 0x35, 0x1A)   # danger red/orange
GREEN     = RGBColor(0x2E, 0x8B, 0x57)   # safe green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x1A, 0x2B, 0x3C)   # body text
GRAY      = RGBColor(0x5B, 0x6B, 0x7B)   # muted text
LIGHT     = RGBColor(0xEE, 0xF2, 0xF7)   # light panel
LIGHT2    = RGBColor(0xE2, 0xE9, 0xF1)   # slightly darker panel
CARD_BORD = RGBColor(0xCF, 0xD9, 0xE4)

FONT_HEAD = "Arial"
FONT_BODY = "Arial"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

# module label shown in footer; updated as we go
_state = {"module": "", "n": 0}


# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------
def _no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, l, t, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    _no_shadow(sp)
    return sp


def grad_rect(s, l, t, w, h, c1, c2, angle=90):
    """Rectangle with a linear two-stop gradient."""
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.line.fill.background()
    _no_shadow(sp)
    # build gradient XML
    spPr = sp.fill._xPr  # noqa
    # remove existing fill
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)
    grad = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = grad.makeelement(qn("a:gsLst"), {})

    def gs(pos, color):
        g = grad.makeelement(qn("a:gs"), {"pos": str(int(pos * 1000))})
        srgb = grad.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
        g.append(srgb)
        return g

    gsLst.append(gs(0, c1))
    gsLst.append(gs(100, c2))
    grad.append(gsLst)
    lin = grad.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000)), "scaled": "1"})
    grad.append(lin)
    # insert gradient before a:ln if present
    ln = spPr.find(qn("a:ln"))
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    return sp


def textbox(s, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size=18, color=INK, bold=False, italic=False, font=FONT_BODY,
         align=PP_ALIGN.LEFT, before=4, after=4, level=0, line=1.05, first=False,
         bullet=None, bullet_color=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(before)
    p.space_after = Pt(after)
    try:
        p.line_spacing = line
    except Exception:
        pass
    r = p.add_run()
    r.text = (bullet + "  " if bullet else "") + text
    f = r.font
    f.size = Pt(size)
    f.name = font
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    if bullet and bullet_color:
        # color just the bullet glyph differently is hard with one run; keep simple
        pass
    return p


# ----------------------------------------------------------------------------
# Shared chrome
# ----------------------------------------------------------------------------
def footer(s, dark=False):
    _state["n"] += 1
    # thin accent line
    rect(s, 0.0, 7.18, 13.333, 0.02, fill=AMBER)
    fg = WHITE if dark else GRAY
    tb, tf = textbox(s, 0.55, 7.16, 8.0, 0.32, anchor=MSO_ANCHOR.MIDDLE)
    para(tf, _state["module"], size=9.5, color=fg, first=True, before=0, after=0)
    tb2, tf2 = textbox(s, 9.5, 7.16, 3.28, 0.32, anchor=MSO_ANCHOR.MIDDLE)
    para(tf2, f"Powerline Electrical Theory & Safety   |   {_state['n']:02d}",
         size=9.5, color=fg, align=PP_ALIGN.RIGHT, first=True, before=0, after=0)


def content_header(s, title, kicker=None):
    # left accent block
    rect(s, 0.0, 0.0, 0.28, 7.5, fill=AMBER)
    rect(s, 0.0, 0.0, 13.333, 0.0)  # noop spacing anchor
    # kicker
    top = 0.45
    if kicker:
        tbk, tfk = textbox(s, 0.7, 0.42, 11.9, 0.34)
        para(tfk, kicker.upper(), size=12.5, color=AMBER_D, bold=True, first=True, before=0, after=0)
        top = 0.78
    tb, tf = textbox(s, 0.68, top, 12.0, 1.0)
    para(tf, title, size=30, color=NAVY, bold=True, first=True, before=0, after=0, line=1.0)
    # underline
    rect(s, 0.7, top + 0.78, 1.7, 0.045, fill=AMBER)


def card(s, l, t, w, h, fill=LIGHT, border=CARD_BORD, line_w=1.0):
    return rect(s, l, t, w, h, fill=fill, line=border, line_w=line_w,
                shape=MSO_SHAPE.ROUNDED_RECTANGLE)


# ----------------------------------------------------------------------------
# Slide templates
# ----------------------------------------------------------------------------
def title_slide(title, subtitle, presenter=None):
    s = slide()
    grad_rect(s, 0, 0, 13.333, 7.5, NAVY, (0x06, 0x16, 0x2B), angle=60)
    # decorative line-tower motif: vertical poles + crossarms
    for x in (10.4, 11.5, 12.6):
        rect(s, x, 0.0, 0.05, 7.5, fill=RGBColor(0x16, 0x33, 0x57))
    for y in (1.6, 2.4, 3.2):
        rect(s, 10.0, y, 3.05, 0.045, fill=RGBColor(0x16, 0x33, 0x57))
    # amber bar
    rect(s, 0.9, 2.25, 0.18, 2.7, fill=AMBER)
    tbk, tfk = textbox(s, 1.3, 2.2, 10.5, 0.5)
    para(tfk, "FIELD TRAINING SERIES", size=15, color=AMBER, bold=True, first=True)
    tb, tf = textbox(s, 1.28, 2.7, 10.6, 2.6)
    for i, line in enumerate(title):
        para(tf, line, size=44, color=WHITE, bold=True, first=(i == 0), before=0, after=2, line=1.0)
    tb2, tf2 = textbox(s, 1.32, 5.15, 10.4, 1.2)
    para(tf2, subtitle, size=18, color=RGBColor(0xC8, 0xD6, 0xE6), first=True, line=1.15)
    if presenter:
        para(tf2, presenter, size=13, color=STEEL, before=12)
    rect(s, 0.0, 7.2, 13.333, 0.06, fill=AMBER)
    _state["n"] += 1
    return s


def section_divider(number, title, items):
    s = slide()
    grad_rect(s, 0, 0, 13.333, 7.5, NAVY, (0x07, 0x18, 0x30), angle=45)
    rect(s, 0.0, 0.0, 13.333, 0.14, fill=AMBER)
    # big module number
    tbn, tfn = textbox(s, 0.85, 1.5, 4.2, 3.2, anchor=MSO_ANCHOR.MIDDLE)
    para(tfn, "MODULE", size=18, color=AMBER, bold=True, first=True, after=0)
    para(tfn, number, size=150, color=WHITE, bold=True, before=0, line=0.9)
    # vertical divider
    rect(s, 5.25, 1.7, 0.04, 4.0, fill=RGBColor(0x2A, 0x47, 0x6E))
    tb, tf = textbox(s, 5.7, 1.85, 7.0, 1.6, anchor=MSO_ANCHOR.TOP)
    para(tf, title, size=34, color=WHITE, bold=True, first=True, line=1.0)
    rect(s, 5.72, 3.35, 1.5, 0.05, fill=AMBER)
    tb2, tf2 = textbox(s, 5.7, 3.65, 7.0, 3.2)
    for i, it in enumerate(items):
        para(tf2, it, size=16, color=RGBColor(0xCF, 0xDC, 0xEC), first=(i == 0),
             before=6, after=6, bullet="—")
    _state["n"] += 1
    return s


def bullets_slide(title, bullets, kicker=None, lead=None):
    """bullets: list of (level, text) or (level, text, bold)"""
    s = slide()
    content_header(s, title, kicker)
    top = 1.95
    if lead:
        tbl, tfl = textbox(s, 0.7, 1.78, 11.9, 0.6)
        para(tfl, lead, size=15.5, color=GRAY, italic=True, first=True)
        top = 2.45
    tb, tf = textbox(s, 0.72, top, 11.95, 7.0 - top - 0.5)
    first = True
    for b in bullets:
        lvl, txt = b[0], b[1]
        bold = b[2] if len(b) > 2 else False
        if lvl == 0:
            para(tf, txt, size=18, color=INK, bold=bold, first=first, level=0,
                 before=7, after=2, bullet="■", line=1.05)
        elif lvl == 1:
            para(tf, txt, size=15, color=RGBColor(0x33, 0x44, 0x55), first=first, level=1,
                 before=2, after=2, bullet="–", line=1.03)
        else:
            para(tf, txt, size=13.5, color=GRAY, first=first, level=2,
                 before=1, after=1, bullet="·", line=1.0)
        first = False
    footer(s)
    return s


def two_col_slide(title, left_head, left_items, right_head, right_items, kicker=None,
                  left_color=NAVY2, right_color=STEEL):
    s = slide()
    content_header(s, title, kicker)
    top = 1.95
    colw = 5.85
    gap = 0.35
    lx = 0.7
    rx = lx + colw + gap
    h = 4.7
    # left card
    card(s, lx, top, colw, h, fill=LIGHT)
    rect(s, lx, top, colw, 0.62, fill=left_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, lx, top + 0.31, colw, 0.31, fill=left_color)  # square off bottom of header
    tbh, tfh = textbox(s, lx + 0.25, top + 0.06, colw - 0.5, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    para(tfh, left_head, size=16, color=WHITE, bold=True, first=True)
    tbl, tfl = textbox(s, lx + 0.3, top + 0.8, colw - 0.6, h - 1.0)
    _fill_items(tfl, left_items)
    # right card
    card(s, rx, top, colw, h, fill=LIGHT)
    rect(s, rx, top, colw, 0.62, fill=right_color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, rx, top + 0.31, colw, 0.31, fill=right_color)
    tbh2, tfh2 = textbox(s, rx + 0.25, top + 0.06, colw - 0.5, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    para(tfh2, right_head, size=16, color=WHITE, bold=True, first=True)
    tbr, tfr = textbox(s, rx + 0.3, top + 0.8, colw - 0.6, h - 1.0)
    _fill_items(tfr, right_items)
    footer(s)
    return s


def _fill_items(tf, items):
    first = True
    for b in items:
        if isinstance(b, str):
            lvl, txt = 0, b
        else:
            lvl, txt = b[0], b[1]
        if lvl == 0:
            para(tf, txt, size=14, color=INK, first=first, bullet="■",
                 before=5, after=2, line=1.03)
        else:
            para(tf, txt, size=12.5, color=GRAY, first=first, level=1, bullet="–",
                 before=1, after=1, line=1.0)
        first = False


def callout_slide(title, body_bullets, callout_kind, callout_title, callout_text, kicker=None):
    """callout_kind: 'danger' | 'warn' | 'safe' | 'info'"""
    s = slide()
    content_header(s, title, kicker)
    top = 1.95
    bodyw = 7.6
    tb, tf = textbox(s, 0.72, top, bodyw, 4.7)
    first = True
    for b in body_bullets:
        lvl, txt = b[0], b[1]
        if lvl == 0:
            para(tf, txt, size=17, color=INK, first=first, bullet="■", before=7, after=2, line=1.05)
        else:
            para(tf, txt, size=14, color=RGBColor(0x33, 0x44, 0x55), first=first, level=1,
                 bullet="–", before=2, after=1, line=1.02)
        first = False
    # callout box on right
    cmap = {
        "danger": (DANGER, RGBColor(0xFB, 0xEA, 0xE7), "DANGER"),
        "warn":   (AMBER_D, RGBColor(0xFD, 0xF4, 0xDE), "WARNING"),
        "safe":   (GREEN, RGBColor(0xE6, 0xF3, 0xEC), "BEST PRACTICE"),
        "info":   (NAVY2, RGBColor(0xE7, 0xEE, 0xF6), "KEY CONCEPT"),
    }
    acc, bg, dlabel = cmap[callout_kind]
    cx, cw = 8.6, 4.13
    ch = 4.3
    cy = top + 0.1
    card(s, cx, cy, cw, ch, fill=bg, border=acc, line_w=1.5)
    rect(s, cx, cy, cw, 0.6, fill=acc, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, cx, cy + 0.3, cw, 0.3, fill=acc)
    tbd, tfd = textbox(s, cx + 0.25, cy + 0.05, cw - 0.5, 0.52, anchor=MSO_ANCHOR.MIDDLE)
    para(tfd, dlabel + "  ⚠" if callout_kind in ("danger", "warn") else dlabel,
         size=14, color=WHITE, bold=True, first=True)
    tbc, tfc = textbox(s, cx + 0.3, cy + 0.78, cw - 0.6, ch - 0.95)
    para(tfc, callout_title, size=15.5, color=acc, bold=True, first=True, after=4)
    if isinstance(callout_text, list):
        for line in callout_text:
            para(tfc, line, size=13, color=INK, bullet="–", before=3, after=1, line=1.05)
    else:
        para(tfc, callout_text, size=13.5, color=INK, line=1.12)
    footer(s)
    return s


def table_slide(title, headers, rows, kicker=None, col_widths=None, note=None, fs=12.5):
    s = slide()
    content_header(s, title, kicker)
    top = 2.0
    ncol = len(headers)
    nrow = len(rows) + 1
    tw = 11.93
    th = min(4.7, 0.5 + 0.46 * len(rows))
    th = max(th, 1.0)
    gtbl = s.shapes.add_table(nrow, ncol, Inches(0.7), Inches(top), Inches(tw), Inches(th)).table
    # disable banding default styling by setting our own
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            gtbl.columns[i].width = Emu(int(Inches(tw) * cw / total))
    # header
    for c, htxt in enumerate(headers):
        cell = gtbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = htxt
        r.font.size = Pt(fs + 0.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT_HEAD
    # body
    for ri, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gtbl.cell(ri, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.03); cell.margin_bottom = Inches(0.03)
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.size = Pt(fs)
            r.font.bold = (c == 0)
            r.font.color.rgb = NAVY if c == 0 else INK
            r.font.name = FONT_BODY
    if note:
        tbn, tfn = textbox(s, 0.72, top + th + 0.18, 11.9, 0.9)
        para(tfn, note, size=12, color=GRAY, italic=True, first=True, line=1.1)
    footer(s)
    return s


def stat_slide(title, stats, kicker=None, lead=None):
    """stats: list of (big, label) -> shown as cards in a row (up to 4)."""
    s = slide()
    content_header(s, title, kicker)
    top = 2.1
    if lead:
        tbl, tfl = textbox(s, 0.7, 1.8, 11.9, 0.55)
        para(tfl, lead, size=15, color=GRAY, italic=True, first=True)
        top = 2.55
    n = len(stats)
    gap = 0.3
    totalw = 11.93
    cw = (totalw - gap * (n - 1)) / n
    h = 2.5
    for i, (big, label) in enumerate(stats):
        x = 0.7 + i * (cw + gap)
        card(s, x, top, cw, h, fill=NAVY)
        rect(s, x, top, cw, 0.12, fill=AMBER, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        tb, tf = textbox(s, x + 0.15, top + 0.35, cw - 0.3, h - 0.5, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, big, size=40, color=AMBER, bold=True, first=True, align=PP_ALIGN.CENTER, after=4, line=0.95)
        para(tf, label, size=13, color=WHITE, align=PP_ALIGN.CENTER, line=1.05)
    # supporting text area below
    tb2, tf2 = textbox(s, 0.72, top + h + 0.3, 11.9, 1.4)
    return s, tf2  # return tf to allow caller to add notes; caller must call footer


def process_slide(title, steps, kicker=None, lead=None):
    """Vertical numbered steps."""
    s = slide()
    content_header(s, title, kicker)
    top = 1.95
    if lead:
        tbl, tfl = textbox(s, 0.7, 1.78, 11.9, 0.55)
        para(tfl, lead, size=15, color=GRAY, italic=True, first=True)
        top = 2.4
    n = len(steps)
    avail = 7.0 - top - 0.5
    rh = min(0.95, avail / n)
    y = top
    for i, (head, body) in enumerate(steps, start=1):
        # number disc
        d = rect(s, 0.7, y, rh - 0.12, rh - 0.12, fill=AMBER, shape=MSO_SHAPE.OVAL)
        tbn, tfn = textbox(s, 0.7, y, rh - 0.12, rh - 0.12, anchor=MSO_ANCHOR.MIDDLE)
        para(tfn, str(i), size=20, color=NAVY, bold=True, first=True, align=PP_ALIGN.CENTER)
        tb, tf = textbox(s, 0.7 + rh + 0.05, y - 0.02, 11.0, rh, anchor=MSO_ANCHOR.MIDDLE)
        para(tf, head, size=15.5, color=NAVY, bold=True, first=True, after=1, line=1.0)
        para(tf, body, size=12.5, color=GRAY, before=0, line=1.0)
        y += rh + 0.06
    footer(s)
    return s


# ============================================================================
# CONTENT
# ============================================================================

# ---- Title ----
title_slide(
    ["Electrical Theory &", "Energized-Line Safety", "for Powerline Workers"],
    "A complete field-training program: theory → working hot → cover-up & tools → job planning & regulations",
    "Lineworker Apprenticeship & Journeyman Refresher  •  Instructor Edition",
)

# ---- Agenda / course map ----
_state["module"] = "Course Overview"
bullets_slide(
    "What This Course Covers",
    [
        (0, "Module 1 — Electrical theory: the science behind everything you touch"),
        (1, "Atomic charge, voltage, current, resistance, power, AC vs DC, the grid"),
        (0, "Module 2 — Electricity and the human body"),
        (1, "Shock physiology, arc flash/blast, step & touch potential"),
        (0, "Module 3 — Working around energized lines & equipment"),
        (1, "Minimum approach distance, the second point of contact, isolation, cover-up, grounding"),
        (0, "Module 4 — Tools of the trade"),
        (1, "Rubber goods, live-line tools, grounds, testers, FR/arc-rated PPE — use, care, testing"),
        (0, "Module 5 — Putting it to work in the field"),
        (1, "Hazard assessment, the tailboard/job brief, roles, communication, rescue"),
        (0, "Module 6 — Regulations & standards that govern the work"),
        (1, "OSHA 1910.269 / 1926 Subpart V, NESC, ASTM/IEEE, recordkeeping"),
    ],
    kicker="Course Map",
    lead="Six modules, building from first principles to the daily disciplines that keep crews alive.",
)

# ---- How to use / objectives ----
two_col_slide(
    "Learning Objectives & How to Use This Training",
    "By the end you will be able to",
    [
        (0, "Explain voltage, current, resistance and power in plain terms"),
        (0, "Describe how the grid moves power from plant to service"),
        (0, "Identify what makes a circuit through the human body"),
        (0, "Recognize and eliminate the second point of contact"),
        (0, "Apply correct cover-up and minimum approach distances"),
        (0, "Build an equipotential grounding zone"),
        (0, "Lead a complete tailboard and hazard assessment"),
        (0, "Cite the regulation behind each safe-work rule"),
    ],
    "How the material is structured",
    [
        (0, "Theory first — you cannot work safely on what you don't understand"),
        (0, "Each hazard is paired with its control"),
        (1, "Recognize → isolate → protect → verify"),
        (0, "DANGER / WARNING / BEST-PRACTICE call-outs flag the must-knows"),
        (0, "Tables summarize numbers you will use on the job"),
        (0, "Discuss every slide as a crew — this is a conversation, not a lecture"),
    ],
    kicker="Objectives",
)

# ---- Why it matters / stats ----
s, tf = stat_slide(
    "Why This Training Exists",
    [("#1", "Electrocution is a leading cause of lineworker fatalities"),
     ("50 V", "Can be enough to be lethal under the right conditions"),
     ("0.1 s", "A fault can evolve to an arc-flash in a fraction of a second"),
     ("100%", "Of electrical contacts are preventable with the right controls")],
    kicker="The Stakes",
    lead="Line work is among the most hazardous trades in the country — and almost every incident traces back to a missed fundamental.",
)
para(tf, "These numbers are not here to scare you — they are here to remind you that the disciplines in this course are written in other workers' experience. Every rule has a reason. The goal of this training is simple: everyone goes home, every day.",
     size=14, color=INK, first=True, line=1.2)
footer(s)

# ============================================================================
# MODULE 1 — ELECTRICAL THEORY
# ============================================================================
_state["module"] = "Module 1 — Electrical Theory"
section_divider("1", "Electrical Theory Fundamentals",
                ["What electricity actually is",
                 "Voltage, current, resistance & power",
                 "AC vs DC and three-phase systems",
                 "Transformers and the grid",
                 "Voltage classes you will encounter"])

bullets_slide(
    "What Is Electricity?",
    [
        (0, "All matter is built from atoms; atoms contain protons (+), neutrons (0) and electrons (–)"),
        (0, "Electricity is the movement of electrons — specifically the flow of electric charge"),
        (1, "In conductors, outer (“free”) electrons drift from atom to atom"),
        (1, "That organized drift, driven by a difference in charge, is electric current"),
        (0, "Two ways to talk about flow direction:"),
        (1, "Conventional current: + to – (the convention used on prints and in formulas)"),
        (1, "Electron flow: – to + (the physical reality)"),
        (0, "Electricity always seeks to return to its source and will take every available path to ground"),
        (1, "It does not “prefer” the easiest path — it uses ALL paths in proportion to their conductance"),
        (1, "Your body can be one of those paths"),
    ],
    kicker="1.1 First Principles",
    lead="Before any rule makes sense, you have to know what is moving in the wire — and why it wants to get to ground.",
)

two_col_slide(
    "Conductors, Insulators & Semiconductors",
    "Conductors (let electrons flow)",
    [
        (0, "Copper, aluminum — line and busbar metals"),
        (0, "Steel — structures, ground rods, messenger"),
        (0, "Water with impurities, wet wood, the human body"),
        (1, "Low resistance — current passes easily"),
        (0, "Why we de-energize, isolate, and ground them"),
    ],
    "Insulators (resist electron flow)",
    [
        (0, "Rubber, dry wood, glass, porcelain, fiberglass"),
        (0, "Air — but only at sufficient distance (it can break down!)"),
        (0, "EPDM rubber — the basis of our gloves & cover-up"),
        (1, "High resistance — used to create separation"),
        (0, "Semiconductors: silicon devices — reclosers, electronics, controls",),
    ],
    kicker="1.2 Materials",
    left_color=DANGER, right_color=GREEN,
)

callout_slide(
    "Voltage — Electrical Pressure",
    [
        (0, "Voltage (V) is the difference in electric potential between two points — the “pressure” that pushes current"),
        (0, "Measured in volts; symbol E (electromotive force) or V"),
        (1, "Analogy: water pressure in a pipe — more pressure, more push"),
        (0, "Voltage is always measured BETWEEN two points (a potential difference)"),
        (1, "A bird on one wire sits at one potential — no difference, no current"),
        (1, "Bridge two different potentials and current flows through you"),
        (0, "Higher voltage can push current across larger air gaps — it can arc to you before you touch it"),
    ],
    "info", "Why this matters",
    ["Every safe-work distance exists because voltage can drive current across air.",
     "“Difference of potential” is the whole game — control it and you control the hazard."],
    kicker="1.3 Voltage",
)

callout_slide(
    "Current — The Flow That Hurts You",
    [
        (0, "Current (I) is the rate of charge flow, measured in amperes (amps)"),
        (1, "Analogy: the volume of water flowing past a point"),
        (0, "It is CURRENT through the body — not voltage — that injures and kills"),
        (1, "Voltage determines how much current will be pushed through your resistance"),
        (0, "Even small currents across the chest can stop the heart"),
        (0, "Fault current on distribution/transmission can be thousands of amps"),
        (1, "Enough to vaporize metal and create an arc-flash explosion"),
    ],
    "danger", "Remember",
    ["It takes surprisingly little current to be fatal.",
     "Milliamps through the heart — not amps — are deadly.",
     "Protect against the current by removing the path and the potential difference."],
    kicker="1.4 Current",
)

bullets_slide(
    "Resistance, Ohm's Law & Power",
    [
        (0, "Resistance (R) opposes current flow — measured in ohms (Ω)"),
        (1, "Conductors = low R; insulators = high R; your skin's R drops sharply when wet"),
        (0, "Ohm's Law ties the three together:  V = I × R"),
        (1, "Current I = V / R   →  same voltage + lower body resistance = MORE current through you"),
        (1, "This is why sweat, water, and broken skin make a shock far more dangerous"),
        (0, "Power is the rate of doing work — measured in watts:  P = V × I"),
        (1, "Also P = I² × R  —  heat in a fault rises with the SQUARE of current"),
        (1, "That heat is what drives burns and arc-flash energy"),
        (0, "Utilities sell energy in kilowatt-hours (kWh): power used over time"),
    ],
    kicker="1.5 Ohm's Law",
    lead="Three quantities, one relationship — the most useful equation you will ever memorize.",
)

table_slide(
    "The Ohm's Law / Power Wheel — Quick Reference",
    ["To find", "If you know", "Use the formula", "Units"],
    [
        ["Voltage (E)", "Current & Resistance", "E = I × R", "Volts (V)"],
        ["Current (I)", "Voltage & Resistance", "I = E ÷ R", "Amperes (A)"],
        ["Resistance (R)", "Voltage & Current", "R = E ÷ I", "Ohms (Ω)"],
        ["Power (P)", "Voltage & Current", "P = E × I", "Watts (W)"],
        ["Power (P)", "Current & Resistance", "P = I² × R", "Watts (W)"],
        ["Power (P)", "Voltage & Resistance", "P = E² ÷ R", "Watts (W)"],
    ],
    kicker="1.5 Reference",
    col_widths=[1.2, 1.6, 1.6, 1.0],
    note="Tip: cover the value you want in the V / I / R triangle and the remaining two show you the operation. Power = the rate energy is delivered — and the source of heat and arc energy.",
)

two_col_slide(
    "AC vs DC",
    "Direct Current (DC)",
    [
        (0, "Electrons flow one direction only"),
        (0, "Batteries, solar arrays, some traction & HVDC systems"),
        (0, "Steady polarity — + and – stay fixed"),
        (0, "Cannot be easily transformed up/down in voltage"),
        (0, "DC can cause sustained muscle 'lock-on'"),
    ],
    "Alternating Current (AC)",
    [
        (0, "Electrons reverse direction cyclically"),
        (0, "60 Hz in North America = 60 cycles per second"),
        (0, "Polarity and magnitude change continuously (a sine wave)"),
        (0, "Easily stepped up/down by transformers — why the grid uses it"),
        (0, "Crosses zero 120×/sec, but is still lethal — never assume 'off'"),
    ],
    kicker="1.6 Waveforms",
    left_color=STEEL, right_color=NAVY2,
)

bullets_slide(
    "AC Details: Frequency, Peak & RMS",
    [
        (0, "One full cycle = one complete positive + negative swing of the sine wave"),
        (0, "Frequency = cycles per second, measured in hertz (Hz); 60 Hz in the U.S./Canada"),
        (0, "Peak voltage is the maximum the wave reaches; the average useful value is the RMS"),
        (1, "RMS (root-mean-square) is the 'DC-equivalent' value — it does the same work"),
        (1, "When we say 120 V or 7,200 V, we mean RMS; peak is ~1.414× higher"),
        (0, "Phase describes where two AC waveforms are in their cycle relative to each other"),
        (1, "Critical for three-phase systems and for paralleling sources"),
        (0, "Inductance and capacitance on lines cause current and voltage to shift out of step",),
        (1, "This is the source of INDUCED voltage on de-energized parallel circuits (Module 3)"),
    ],
    kicker="1.6 AC Behavior",
)

two_col_slide(
    "Single-Phase vs Three-Phase",
    "Single-Phase",
    [
        (0, "One AC voltage waveform (with neutral)"),
        (0, "Typical residential service: 120/240 V"),
        (0, "Simple loads — lighting, receptacles, small motors"),
        (0, "Power delivery pulses (crosses zero each half-cycle)"),
    ],
    "Three-Phase",
    [
        (0, "Three waveforms 120° apart — A, B, C phases"),
        (0, "Backbone of transmission, distribution & industrial power"),
        (0, "Smooth, continuous power — efficient for large motors"),
        (0, "More power in less conductor than single-phase"),
        (0, "Phase-to-phase voltage = √3 (1.732) × phase-to-neutral"),
    ],
    kicker="1.7 Phases",
    left_color=STEEL, right_color=NAVY2,
)

bullets_slide(
    "Transformers — Stepping Voltage Up & Down",
    [
        (0, "A transformer changes AC voltage using two coils linked by a magnetic core"),
        (1, "No direct electrical connection — energy crosses by induction"),
        (0, "Turns ratio sets the voltage ratio:  more secondary turns → higher voltage"),
        (1, "Step UP: raise voltage for transmission (less current = less line loss)"),
        (1, "Step DOWN: lower voltage for distribution and the customer"),
        (0, "Power in ≈ power out, so when voltage goes UP, current goes DOWN (and vice-versa)"),
        (1, "P = V × I — this is WHY we transmit at high voltage: to cut I²R losses"),
        (0, "Distribution transformers on the pole/pad drop primary (e.g., 7,200 V) to 120/240 V",),
        (0, "WARNING: a de-energized secondary can be BACK-FED from the customer side — generators!"),
    ],
    kicker="1.8 Transformers",
    lead="The single reason the grid runs on AC — voltage can be transformed efficiently.",
)

process_slide(
    "The Grid: From Generation to Your Service",
    [
        ("1. Generation", "Power plants (gas, hydro, nuclear, wind, solar) produce AC, typically at medium voltage."),
        ("2. Step-Up Substation", "Transformers raise voltage to transmission levels (69 kV up to 765 kV) to move power far with low loss."),
        ("3. Transmission Lines", "High-voltage lines on large towers carry bulk power across regions."),
        ("4. Step-Down / Distribution Substation", "Voltage is reduced to distribution levels (typically 4 kV–35 kV)."),
        ("5. Distribution Feeders", "Primary lines on poles carry power through neighborhoods; reclosers & fuses protect them."),
        ("6. Service Transformer & Customer", "Final step-down to 120/240 V (residential) or higher for commercial/industrial."),
    ],
    kicker="1.9 The System",
    lead="Know where you are on this chain — it tells you the voltage, the fault energy, and the controls you need.",
)

table_slide(
    "Voltage Classes You Will Work Around",
    ["Class", "Typical Range", "Where you see it", "Note"],
    [
        ["Low voltage", "< 600 V", "Services, secondary, street light", "Still deadly — most contacts happen here"],
        ["Medium voltage (Distribution)", "600 V – 35 kV", "Primary feeders, urban/rural lines", "Common lineworker environment"],
        ["High voltage (Sub-transmission/Transmission)", "35 kV – 230 kV", "Towers, large substations", "Larger MAD; arc & induction hazards grow"],
        ["Extra-high voltage (EHV)", "230 kV – 765 kV", "Bulk transmission", "Specialized live-line methods, large clearances"],
        ["Ultra-high voltage (UHV)", "> 765 kV", "Long-distance bulk corridors", "Rare; highly specialized work"],
    ],
    kicker="1.10 Voltage Classes",
    col_widths=[1.5, 1.2, 1.7, 1.7],
    note="Exact class boundaries vary by utility and standard (NESC/IEEE). What never changes: higher voltage = larger required clearances and greater fault energy.",
)

# ----------------------------------------------------------------------------
# Embedded illustration slides (rendered from diagrams.js -> assets/*.png)
# ----------------------------------------------------------------------------
def image_slide(title, fig_key, caption):
    s = slide()
    content_header(s, title, kicker="Illustrated")
    png = os.path.join(ASSETS, "fig_%s.png" % fig_key)
    iw, ih = Image.open(png).size
    box_l, box_t, box_w, box_h = 0.7, 1.95, 11.93, 4.05
    ar = iw / ih
    w = box_w
    h = w / ar
    if h > box_h:
        h = box_h
        w = h * ar
    left = box_l + (box_w - w) / 2.0
    top = box_t + (box_h - h) / 2.0
    s.shapes.add_picture(png, Inches(left), Inches(top), width=Inches(w), height=Inches(h))
    tb, tf = textbox(s, 0.9, 6.18, 11.5, 0.7, anchor=MSO_ANCHOR.TOP)
    para(tf, caption, size=13, color=GRAY, italic=True, first=True, align=PP_ALIGN.CENTER, line=1.18)
    footer(s)
    s.notes_slide.notes_text_frame.text = re.sub(r"<[^>]+>", "", caption)
    return s


MODULE_FIGS = {
    1: [
        ("atom", "How Electricity Flows", "Electrons drifting from atom to atom — driven by a difference in charge — are electric current."),
        ("conductor", "Conductors vs Insulators", "Conductors let electrons pass (low resistance); insulators bind them (high resistance)."),
        ("water", "The Water Analogy", "Voltage is pressure, current is flow, resistance is the restriction:  I = V ÷ R."),
        ("bird", "Why Voltage Needs Two Points", "One potential is safe; bridging two different potentials drives current through you."),
        ("ohm", "The Ohm's Law & Power Wheel", "Cover the unknown in the V / I / R triangle. Power P = V×I = I²R — heat rises with current squared."),
        ("acdc", "AC vs DC Waveforms", "DC flows one way at a steady level; AC reverses 60×/sec, which lets transformers change its voltage."),
        ("threephase", "Three-Phase Power", "Three waveforms 120° apart deliver smooth, continuous power — the backbone of the grid."),
        ("transformer", "How a Transformer Works", "Two coils on a magnetic core: the turns ratio sets the voltage ratio; power in ≈ power out."),
        ("grid", "The Grid, End to End", "Voltage is stepped up for transmission, then down in stages until it is safe for the customer."),
    ],
    2: [
        ("body", "Current Path Through the Body", "Current across the chest — hand-to-hand or hand-to-foot — can stop the heart. Control both contact points."),
        ("threshold", "Shock Threshold Ladder", "The dangerous range begins in milliamps — far below what trips a typical breaker."),
        ("arcflash", "Arc Flash & Arc Blast", "An arc can reach ~35,000°F with an explosive blast — and can be triggered by proximity alone."),
        ("steptouch", "Step & Touch Potential", "Current spreading from a downed line creates a ground gradient. Shuffle out with feet together."),
        ("recloser", "Reclosers — Why a Dead Line Bites Back", "Reclosers trip then auto-reclose — a downed line can re-energize without warning."),
    ],
    3: [
        ("hierarchy", "The Hierarchy of Controls", "Work from the top down: de-energize first; PPE is the last line, never the only line."),
        ("mad", "Minimum Approach Distance", "Keep body, tools, and equipment outside the minimum approach distance for your voltage."),
        ("bucket", "The Insulated Aerial Device", "The fiberglass upper boom is a dielectric insulating gap between the bucket and ground."),
        ("secondpoint", "The Second Point of Contact", "Eliminate the second point of contact — cover it or remove it — and the circuit can't close through you."),
        ("coverup", "Cover-Up at the Pole Top", "Insulate every energized and grounded part within reach: nearest first, remove in reverse."),
        ("loto", "The De-Energize / LOTO Sequence", "Identify → open → lock & tag → test → ground.  Never skip a step."),
        ("epz", "The Equipotential Zone", "Bond and ground a work zone to one potential — no difference across the body means no current."),
        ("induced", "How Induced Voltage Appears", "A parallel live line couples voltage onto a dead one — ground at the work site, not just the ends."),
    ],
    4: [
        ("gloves", "Rubber Glove Voltage Classes", "Select a glove class rated at or above your system voltage; wear leather protectors over them."),
        ("hotstick", "The Live-Line (Hot-Stick) Method", "A fiberglass live-line tool extends reach beyond MAD — keep it clean and dry."),
        ("phasing", "Phasing & Voltage Testing", "Verify voltage and phase — and prove the tester on a known source before and after."),
        ("grounds", "Protective Grounding Sets", "Grounds carry fault current and trip protection fast: ground end first, line ends last."),
        ("ppe", "Head-to-Toe PPE", "Arc-rated FR clothing, Class E hard hat, face shield, rubber gloves — the last layer of defense."),
    ],
    5: [
        ("tailboard", "The Tailboard Briefing", "The crew briefing aligns everyone: hazards, procedures, controls, PPE, roles, and rescue."),
        ("switching", "Switching & the Clearance", "A clearance isolates the work zone: open, locked, tagged, tested dead, and grounded on every side."),
        ("rescue", "Pole-Top Rescue", "Remove the source first, lower the patient, start CPR & AED — never become the second victim."),
    ],
    7: [
        ("arcenergy", "Incident Energy vs Working Distance", "Energy falls with the square of distance; where it crosses 1.2 cal/cm² is the arc-flash boundary."),
        ("ppecat", "Arc-Flash PPE Categories", "Four categories from ≥4 to ≥40 cal/cm². The garment's arc rating must meet or exceed the incident energy."),
        ("arclabel", "The Arc-Flash & Shock Label", "Voltage, boundaries, incident energy, PPE category, and glove class — read it before you open the door."),
        ("scenecontrol", "Emergency Scene Control", "Treat the conductor and ground as energized; keep the patient untouched until the source is removed — then reach them."),
        ("cpraed", "CPR & AED Sequence", "Once the scene is safe: check & call, compress 100–120/min, apply the AED, continue until EMS arrives."),
    ],
}


def add_module_figs(n):
    for fig, title, cap in MODULE_FIGS.get(n, []):
        image_slide(title, fig, cap)


add_module_figs(1)
bullets_slide(
    "Module 1 — Key Takeaways",
    [
        (0, "Electricity is the flow of charge that always wants to return to its source via every path"),
        (0, "Voltage pushes; current flows; resistance opposes — V = I × R ties them together"),
        (0, "It is CURRENT through the body that kills, and low body resistance lets more through"),
        (0, "AC dominates the grid because transformers let us change voltage efficiently"),
        (0, "Transmit high (low current, low loss); distribute and use low (safe for customers)"),
        (0, "Know your voltage class — it sets your clearances, fault energy, and controls"),
    ],
    kicker="Module 1 Recap",
)

# ============================================================================
# MODULE 2 — ELECTRICITY & THE HUMAN BODY
# ============================================================================
_state["module"] = "Module 2 — Electricity & the Body"
section_divider("2", "Electricity & the Human Body",
                ["How current injures and kills",
                 "Shock thresholds & the heart",
                 "Electrical burns",
                 "Arc flash and arc blast",
                 "Step, touch & ground-gradient potential"])

bullets_slide(
    "How Current Affects the Body",
    [
        (0, "The body is a conductor — nerves and muscles run on tiny electrical signals"),
        (0, "External current overrides those signals: muscles clamp, nerves misfire, the heart can fibrillate"),
        (0, "Severity depends on:"),
        (1, "Amount of current (the key factor) — set by voltage ÷ your resistance"),
        (1, "Path through the body — hand-to-hand or hand-to-foot crosses the heart"),
        (1, "Duration of contact — longer exposure, more damage"),
        (1, "Frequency — 60 Hz AC is especially effective at causing fibrillation"),
        (0, "Skin resistance can be ~100,000 Ω dry, but drops to ~1,000 Ω or less when wet/broken"),
        (1, "Same voltage then drives far more current — wet conditions are a force multiplier"),
    ],
    kicker="2.1 Shock Physiology",
)

table_slide(
    "Approximate Effects of 60 Hz Current Across the Body",
    ["Current (through body)", "Typical effect"],
    [
        ["~1 mA", "Threshold of perception — a faint tingle"],
        ["~5 mA", "Slight shock; not painful but startling — can cause falls"],
        ["6–16 mA", "Painful shock; loss of muscular control begins (“let-go” threshold)"],
        ["17–20 mA", "Muscles clamp — CANNOT let go of the conductor; breathing labored"],
        ["20–50 mA", "Severe muscle contractions, breathing difficulty"],
        ["50–100 mA", "Possible ventricular fibrillation — the heart quivers, stops pumping"],
        ["100 mA – 2 A", "Likely fatal ventricular fibrillation; nerve damage"],
        ["> 2 A", "Cardiac arrest, severe internal burns, organ damage"],
    ],
    kicker="2.2 Thresholds",
    col_widths=[1.3, 3.2],
    note="Values are approximate and vary by body size, path, and conditions. The lesson: the dangerous range starts in the MILLIAMP region — far below what trips most circuits.",
)

callout_slide(
    "The Heart, Let-Go & Fibrillation",
    [
        (0, "“Let-go threshold”: above ~10–16 mA your flexor muscles clamp and you cannot release the conductor"),
        (1, "The longer you hold on, the more the contact heats and the more current flows"),
        (0, "Ventricular fibrillation: current disrupts the heart's rhythm so it quivers instead of pumping"),
        (1, "Blood stops circulating — brain damage in minutes, death without intervention"),
        (1, "A defibrillator/AED and CPR may restart a coordinated rhythm — know where yours is"),
        (0, "Even a 'minor' shock can throw you from a pole or bucket — the fall often does the damage"),
    ],
    "danger", "Why we never rely on 'holding on'",
    ["Above the let-go threshold your body works against you.",
     "The only safe contact with an energized conductor is NO direct contact — insulate and isolate."],
    kicker="2.3 The Heart",
)

two_col_slide(
    "Electrical Burns",
    "How they happen",
    [
        (0, "Contact / true electrical burns: current flowing through tissue generates heat (I²R)"),
        (0, "Arc burns: radiated heat from an arc — can exceed 35,000°F at the arc"),
        (0, "Thermal-contact burns: ignited clothing or hot surfaces"),
        (0, "Flash burns to eyes (welder's flash) from intense UV"),
    ],
    "Why they are deceptive",
    [
        (0, "Entry/exit wounds may look small but mask massive internal damage"),
        (0, "Current cooks tissue along its whole path — nerves, muscle, organs"),
        (0, "Delayed effects: kidney damage from tissue breakdown, cardiac issues"),
        (0, "ALWAYS seek medical evaluation after any electrical contact — even if you feel fine"),
    ],
    kicker="2.4 Burns",
    left_color=DANGER, right_color=AMBER_D,
)

callout_slide(
    "Arc Flash & Arc Blast",
    [
        (0, "An arc flash is an explosive release of energy when current jumps through ionized air"),
        (1, "Triggered by faults, dropped tools, animals, equipment failure, or getting too close"),
        (0, "Arc FLASH = light & heat: temperatures up to ~35,000°F (hotter than the sun's surface)"),
        (0, "Arc BLAST = pressure wave: vaporized metal expands violently"),
        (1, "Copper expands ~67,000 times when it vaporizes — like a small explosion"),
        (1, "Concussive force, shrapnel, molten metal spray, sound > 140 dB"),
        (0, "Incident energy is measured in cal/cm² — it sets the arc rating of PPE you must wear"),
    ],
    "danger", "Controlled by",
    ["Distance (clearances), de-energizing, and arc-rated FR clothing/face protection.",
     "An arc flash can occur WITHOUT you touching anything — proximity alone can trigger it."],
    kicker="2.5 Arc Flash",
)

callout_slide(
    "Step Potential & Touch Potential",
    [
        (0, "When current flows into the earth (e.g., a downed line, a fault to a ground rod), it spreads out"),
        (0, "Voltage drops with distance from the contact point — forming a 'voltage gradient' across the ground"),
        (0, "STEP POTENTIAL: the voltage difference between your two feet"),
        (1, "One foot sits at a higher potential than the other — current flows up one leg, down the other"),
        (1, "The wider your stride, the greater the difference — and the greater the shock"),
        (0, "TOUCH POTENTIAL: difference between your hand (on an object) and your feet"),
        (1, "Touching a faulted structure while standing on ground at a different potential"),
    ],
    "danger", "If you must leave an energized area",
    ["Keep feet together and SHUFFLE or bunny-hop — never take a normal stride.",
     "Assume the ground around any downed line is energized out to a wide radius.",
     "Equipotential bonding & grounding (Module 3) eliminate these gradients in your work zone."],
    kicker="2.6 Ground Gradients",
)

bullets_slide(
    "Downed Conductors & The Ground Around You",
    [
        (0, "A downed line may look dead — it is NOT until tested and grounded"),
        (1, "It may be re-energized automatically by reclosers, or back-fed by a customer's generator"),
        (0, "Energized ground can extend many feet from the contact point — you cannot see it"),
        (0, "Vehicles in contact with lines: occupants should STAY INSIDE unless fire forces evacuation"),
        (1, "If they must exit: JUMP clear without touching vehicle and ground at once, then shuffle away"),
        (0, "Keep the public back — establish a wide barrier and treat everything as energized"),
        (0, "Never approach within the established radius until the line is confirmed de-energized & grounded"),
    ],
    kicker="2.7 Downed Lines",
    lead="The most dangerous ground is the ground that looks perfectly normal.",
)

add_module_figs(2)
bullets_slide(
    "Module 2 — Key Takeaways",
    [
        (0, "Current — not voltage — injures; the dangerous range begins in milliamps"),
        (0, "Path across the heart and 60 Hz AC make fibrillation likely; wet skin worsens everything"),
        (0, "Above the let-go threshold you cannot release — so we never rely on letting go"),
        (0, "Arc flash/blast can kill or maim WITHOUT contact — distance and FR PPE are your defense"),
        (0, "Energized ground creates step & touch potential — shuffle out, bond and ground your zone"),
        (0, "After any contact: get medical evaluation — internal damage hides"),
    ],
    kicker="Module 2 Recap",
)

# ============================================================================
# MODULE 3 — WORKING AROUND ENERGIZED LINES
# ============================================================================
_state["module"] = "Module 3 — Working Energized"
section_divider("3", "Working Around Energized Lines & Equipment",
                ["Minimum approach distance",
                 "The second point of contact — the core concept",
                 "Isolate, insulate & cover-up",
                 "Grounding & the equipotential zone",
                 "De-energizing, LOTO & induced voltage"])

bullets_slide(
    "The Hierarchy of Controls for Electrical Work",
    [
        (0, "The safest energized work is work that isn't energized — always prefer de-energizing where feasible"),
        (0, "Order of preference when planning a job:"),
        (1, "1. Eliminate — de-energize, isolate, lock & tag, test dead, and ground"),
        (1, "2. Substitute / engineer — reroute, switch load, use remote/robotic methods"),
        (1, "3. Isolate the worker — insulate & isolate (rubber-glove or hot-stick / live-line bare-hand)"),
        (1, "4. Administrative — procedures, job briefings, MAD, qualified workers only"),
        (1, "5. PPE — FR clothing, rubber gloves/sleeves, face & head protection (last line, never the only line)"),
        (0, "Energized work is permitted only when de-energizing is infeasible or creates greater hazard"),
    ],
    kicker="3.1 Controls",
    lead="Reach for the top of the hierarchy first; PPE alone is never the plan.",
)

callout_slide(
    "Minimum Approach Distance (MAD)",
    [
        (0, "MAD = the closest a qualified worker (or conductive object) may come to an energized part"),
        (0, "It protects against flashover — the air gap that voltage cannot jump"),
        (0, "MAD increases with voltage and with altitude (thinner air insulates less)"),
        (0, "It includes the phase-to-ground AND phase-to-phase exposure"),
        (0, "Anything conductive counts: your body, tools, materials, the boom, a tape measure"),
        (1, "Unqualified employees must stay even farther back"),
        (0, "If you cannot maintain MAD, the part must be de-energized OR covered/insulated"),
    ],
    "warn", "Rule of thumb",
    ["Know your utility's MAD table for the voltage you're on — BEFORE you leave the ground.",
     "When in doubt, cover up or back off. MAD is a minimum, not a target."],
    kicker="3.2 MAD",
)

table_slide(
    "Illustrative Minimum Approach Distances (Phase-to-Ground)",
    ["Nominal Voltage (ph-ph)", "Example MAD", "Typical Setting"],
    [
        ["50 V – 300 V", "Avoid contact", "Service / secondary"],
        ["301 V – 750 V", "~1 ft 1 in", "Secondary / street light"],
        ["751 V – 15 kV", "~2 ft 2 in", "Distribution primary"],
        ["15.1 kV – 36 kV", "~2 ft 7 in", "Distribution / sub-transmission"],
        ["36.1 kV – 46 kV", "~2 ft 9 in", "Sub-transmission"],
        ["46.1 kV – 72.5 kV", "~3 ft 3 in", "Sub-transmission"],
        ["72.6 kV – 121 kV", "~3 ft 4 in", "Transmission"],
        ["138 kV – 145 kV", "~3 ft 7 in", "Transmission"],
        ["230 kV", "~5 ft 3 in", "Bulk transmission"],
    ],
    kicker="3.2 MAD Table",
    col_widths=[1.6, 1.2, 1.6],
    fs=11.5,
    note="ILLUSTRATIVE ONLY — distances vary with altitude, transient over-voltage factor, and the governing edition of OSHA 1910.269 Table R-3 / R-6. ALWAYS use your employer's current, approved MAD table.",
)

# THE CORE CONCEPT
callout_slide(
    "The Second Point of Contact — The Concept That Saves Lives",
    [
        (0, "Current can only flow when it has a complete path: a source, a load/body, and a return"),
        (0, "Touching ONE energized conductor while perfectly insulated from everything else = no path"),
        (1, "This is why a bird sits on a wire unharmed — it is at one potential, with no second contact"),
        (0, "The SECOND POINT OF CONTACT is what completes the circuit through you:"),
        (1, "A grounded neutral, a guy wire, a pole ground, a crossarm, the structure"),
        (1, "Another phase at a different potential (phase-to-phase — even more dangerous)"),
        (1, "A co-worker, a tool, the bucket, a wet rope, the truck"),
        (0, "Eliminate the second point of contact and you eliminate the shock — every time"),
    ],
    "danger", "The discipline",
    ["Before you reach for ANY energized part, ask: 'Where is my second point of contact?'",
     "Cover it. Insulate it. Or remove yourself from it.",
     "Most electrocutions are NOT phase-to-phase — they are phase-to-GROUND through a missed second point."],
    kicker="3.3 SECOND POINT OF CONTACT",
)

bullets_slide(
    "Finding Your Second Point of Contact — Every Time",
    [
        (0, "Scan the whole work envelope, not just the conductor you're working"),
        (0, "Common second points hiding in plain sight:"),
        (1, "Grounded neutral conductor (often the closest deadly path)"),
        (1, "The pole ground / ground wire, guy wires and down-guys"),
        (1, "Metal crossarm braces, bolts, staples, hardware"),
        (1, "Other phases above, below, or behind you"),
        (1, "Transformer cases, capacitor banks, riser/cable terminations"),
        (1, "The bucket lip, jib, winch line, hand line, or another worker"),
        (0, "Work to keep yourself, your tools, and the energized part isolated from all of them"),
        (0, "Glove AND cover: rubber gloves protect your hands; cover-up protects against the SECOND point"),
    ],
    kicker="3.3 Application",
    lead="Rubber gloves alone are not the answer — you must also control what else you could touch.",
)

bullets_slide(
    "Insulate & Isolate — The Two Methods of Protection",
    [
        (0, "INSULATE: place insulating barriers between the worker and energized/grounded parts"),
        (1, "Rubber gloves, sleeves, line hose, hoods, blankets — rated for the voltage"),
        (1, "Insulated tools and platforms; the dielectric boom of an aerial device"),
        (0, "ISOLATE: create distance/separation so no path can form"),
        (1, "Maintain MAD; reposition the work; use live-line (hot-stick) tools to keep the worker back"),
        (0, "Rubber-glove method: worker is insulated and works within reach — cover-up is mandatory"),
        (0, "Hot-stick / live-line tool method: worker stays beyond MAD and works through tools"),
        (0, "Bare-hand (live-line) method: worker is bonded TO the conductor at line potential — specialized, EHV"),
    ],
    kicker="3.4 Two Methods",
)

bullets_slide(
    "Cover-Up: Purpose & Principles",
    [
        (0, "Cover-up = installing insulating guards over energized AND grounded parts in the work zone"),
        (0, "Purpose: protect against accidental contact and eliminate second points of contact"),
        (1, "It is secondary protection — it backs up your gloves, MAD, and positioning"),
        (0, "Cover EVERYTHING within reach you could contact: phases, neutral, ground, hardware, guys"),
        (0, "Cover-up must be rated for the system voltage and inspected before use"),
        (0, "Install from a position of safety — nearest/most-accessible energized part first, working outward",),
        (1, "Use a hot stick or gloved hands per procedure; never let cover-up create a false sense of security"),
        (0, "Remove in REVERSE order — last on, first off — keeping protection in place longest where you work"),
    ],
    kicker="3.5 Cover-Up",
    lead="Cover-up turns a cluttered, energized pole top into an isolated, defendable work zone.",
)

process_slide(
    "Cover-Up Sequence — Building the Protected Zone",
    [
        ("1. Plan & inspect", "Identify every energized & grounded part within reach. Inspect all cover-up for damage; verify voltage rating."),
        ("2. Establish footing", "Position the bucket/yourself so you can reach with the conductor at a safe angle — from below where possible."),
        ("3. Cover the closest hazard first", "Apply line hose/blankets to the nearest energized phase, then work outward to additional phases."),
        ("4. Cover the second points", "Cover the neutral, ground wire, guys, crossarm hardware — anything that could be your second point of contact."),
        ("5. Work the job", "With the zone insulated, perform the task maintaining gloves, MAD, and positioning throughout."),
        ("6. Remove in reverse", "Strip cover-up last-on/first-off — keep protection nearest you in place until the end."),
    ],
    kicker="3.5 Sequence",
)

two_col_slide(
    "Cover-Up Equipment",
    "Types of cover-up",
    [
        (0, "Line hose — tubular cover for conductors"),
        (0, "Insulator hoods — cover pin/post insulators & tops"),
        (0, "Conductor/line covers & connector covers"),
        (0, "Insulating blankets — wrap irregular shapes, hardware"),
        (0, "Dead-end & pole/crossarm covers, guy guards"),
    ],
    "Selection & care",
    [
        (0, "Match the voltage rating to (or above) the system"),
        (0, "Inspect for cuts, ozone cracks, embedded metal, swelling"),
        (0, "Keep clean & dry; store properly — no sharp folds"),
        (0, "Use approved clip pins/connectors so it stays put"),
        (0, "Removed from service & retested per schedule if suspect"),
    ],
    kicker="3.5 Equipment",
    left_color=NAVY2, right_color=GREEN,
)

callout_slide(
    "De-Energizing & Lockout/Tagout (LOTO)",
    [
        (0, "When a line can be de-energized, follow the formal switching & clearance procedure"),
        (0, "The classic verification sequence — do not skip a step:"),
        (1, "1. Identify the correct circuit/section (right line, right structure)"),
        (1, "2. Open the disconnects and render them visibly open"),
        (1, "3. Lock and TAG the points of isolation — each worker's own lock"),
        (1, "4. TEST for absence of voltage with a tested-good detector (test-before-touch)"),
        (1, "5. Apply protective GROUNDS to create an equipotential, drained zone"),
        (0, "Hold the clearance; only the holder releases it after all workers are clear"),
    ],
    "warn", "Test your tester",
    ["Prove the voltage detector works on a known source BEFORE and AFTER testing the line.",
     "Treat every conductor as energized until YOU have verified it dead and grounded."],
    kicker="3.6 De-Energize / LOTO",
)

callout_slide(
    "Grounding & The Equipotential Zone",
    [
        (0, "Protective grounding does NOT just 'drain' the line — it protects YOU if it is re-energized"),
        (0, "If the line is accidentally energized, grounds provide a low-resistance path that trips protection fast"),
        (0, "Equipotential zone (EPZ): bond everything you could touch to the SAME potential"),
        (1, "If everything around you is at the same voltage, no current flows through you — like the bird"),
        (0, "Bonding: connect conductors, structures, and equipment together; ground that bonded system"),
        (0, "Grounds must be sized for available fault current and connected in the correct order"),
        (1, "Connect to ground/structure FIRST, then to the conductor; remove in REVERSE"),
    ],
    "info", "Why grounds save lives",
    ["A worker inside a properly bonded EPZ rides the voltage up and down together — no difference of potential across the body.",
     "Grounds are sized to carry fault current long enough for the breaker/recloser to operate."],
    kicker="3.7 Grounding & EPZ",
)

bullets_slide(
    "Induced Voltage — The Hidden Hazard on 'Dead' Lines",
    [
        (0, "A de-energized line running parallel to an energized one can pick up voltage by induction"),
        (1, "Magnetic (current) and electrostatic (voltage) coupling from the live circuit"),
        (0, "The longer the parallel and the higher the live circuit's load, the greater the induced voltage"),
        (0, "An ungrounded 'dead' line can hold enough induced potential to shock or kill"),
        (0, "Control it by GROUNDING — and by working within a bonded equipotential zone"),
        (1, "Grounds at the work site (not just at the ends) keep induced potential off the worker"),
        (0, "Switching, transmission, and parallel-circuit work demand extra attention to induction"),
    ],
    kicker="3.8 Induced Voltage",
    lead="“Dead” is not the same as “grounded.” Induction is why we ground at the work location.",
)

add_module_figs(3)
bullets_slide(
    "Module 3 — Key Takeaways",
    [
        (0, "Prefer de-energizing; when working hot, isolate AND insulate — never rely on one control"),
        (0, "Maintain MAD for your voltage; anything conductive counts toward the distance"),
        (0, "THE SECOND POINT OF CONTACT completes the circuit — find it, cover it, remove it, every time"),
        (0, "Cover-up backs up your gloves and MAD: nearest first, remove in reverse"),
        (0, "De-energize properly: identify, open, lock/tag, TEST, then GROUND"),
        (0, "Bond an equipotential zone and ground at the work site — it defeats re-energization AND induction"),
    ],
    kicker="Module 3 Recap",
)

# ============================================================================
# MODULE 4 — TOOLS OF THE TRADE
# ============================================================================
_state["module"] = "Module 4 — Tools of the Trade"
section_divider("4", "Tools of the Trade",
                ["Rubber insulating gloves & sleeves",
                 "Live-line (hot-stick) tools",
                 "Protective grounds & testers",
                 "Arc-rated FR clothing & PPE",
                 "Inspection, testing & care"])

bullets_slide(
    "Rubber Insulating Gloves — Your First Line of Defense",
    [
        (0, "Rubber gloves insulate your hands from energized parts in the rubber-glove method"),
        (0, "Always worn with leather protector gloves OVER them"),
        (1, "The leather shields the rubber from cuts, abrasion, and UV — the rubber does the insulating"),
        (0, "Selected by VOLTAGE CLASS — each class has a max use voltage and proof-test voltage"),
        (0, "Field-inspect before EVERY use:"),
        (1, "Air-test (roll/inflate) to find pinholes; check for cuts, cracks, swelling, embedded debris"),
        (1, "Inspect inside and out; check for chemical or ozone damage"),
        (0, "Retest electrically on a schedule (typically every 6 months) at a certified lab"),
        (0, "Store in a glove bag, fingers up, away from heat, ozone, sunlight, and sharp objects"),
    ],
    kicker="4.1 Rubber Gloves",
)

table_slide(
    "Rubber Glove Voltage Classes (ASTM D120 / IEC)",
    ["Class", "Max Use Voltage (AC, ph-ph)", "Color (label)", "Proof Test"],
    [
        ["Class 00", "500 V", "Beige", "2,500 V"],
        ["Class 0", "1,000 V", "Red", "5,000 V"],
        ["Class 1", "7,500 V", "White", "10,000 V"],
        ["Class 2", "17,000 V", "Yellow", "20,000 V"],
        ["Class 3", "26,500 V", "Green", "30,000 V"],
        ["Class 4", "36,000 V", "Orange", "40,000 V"],
    ],
    kicker="4.1 Glove Classes",
    col_widths=[0.8, 1.8, 1.1, 1.1],
    note="Always select a class rated AT or ABOVE the maximum voltage you may contact. Label color aids quick identification; the printed rating governs. Same class system applies to insulating sleeves and many blankets.",
)

two_col_slide(
    "Rubber Sleeves, Blankets & Line Hose",
    "Sleeves",
    [
        (0, "Insulate the forearm, upper arm, and shoulder"),
        (0, "Worn when the arm could contact an energized/grounded part"),
        (0, "Same voltage-class system as gloves"),
        (0, "Inspected & retested like gloves"),
    ],
    "Blankets & line hose",
    [
        (0, "Blankets: flexible cover for irregular shapes & hardware"),
        (0, "Line hose: slips over conductors to insulate runs"),
        (0, "Secured with approved pins/snaps/clamps"),
        (0, "Voltage-rated; inspected for cuts, ozone & contamination"),
    ],
    kicker="4.2 Sleeves & Cover",
    left_color=NAVY2, right_color=STEEL,
)

bullets_slide(
    "Live-Line (Hot-Stick) Tools",
    [
        (0, "Hot sticks are insulating tools (usually fiberglass/epoxy) that extend your reach beyond MAD"),
        (0, "Let qualified workers operate, test, install and remove equipment without closing the gap"),
        (0, "Common types:"),
        (1, "Universal / shotgun (telescoping) sticks with interchangeable heads"),
        (1, "Switch sticks, disconnect sticks, fuse pullers"),
        (1, "Tie sticks, wire-holding & cleaning tools, measuring sticks"),
        (1, "Grounding sticks (clamp sticks) for applying protective grounds"),
        (0, "Keep them CLEAN and DRY — contamination & moisture ruin the insulating surface"),
        (1, "Wipe down with an approved cleaner/wax; use a cradle, not the ground"),
        (0, "Inspect for scratches, cracks, soft spots; periodic dielectric testing required"),
    ],
    kicker="4.3 Hot Sticks",
    lead="A hot stick keeps the worker back while the work goes forward — isolation in tool form.",
)

bullets_slide(
    "Protective Grounds & Grounding Equipment",
    [
        (0, "Personal protective grounds = flexible copper cables with clamps (ferrules) and ground studs"),
        (0, "Sized to carry the available fault current for the clearing time without failing"),
        (1, "Undersized grounds can fuse and fail under fault — use the correct cable size"),
        (0, "Components: cable, clamps, ferrules, ground rods/cluster bars, EPZ jumpers"),
        (0, "Apply correctly: clean contact points, tight clamps, ground end FIRST then line end"),
        (0, "Use grounding sticks to install/remove — treat the line as energized until grounds are on"),
        (0, "Inspect for broken strands, corroded clamps, loose ferrules; remove damaged sets from service"),
        (0, "Remove grounds in REVERSE: line end first, ground end last"),
    ],
    kicker="4.4 Grounds",
)

two_col_slide(
    "Voltage Detectors & Testers",
    "Types",
    [
        (0, "Contact voltage detectors (on a hot stick)"),
        (0, "Non-contact / proximity testers (presence of field)"),
        (0, "Phasing testers — verify phase relationships"),
        (0, "Audible/visual indication of energized condition"),
    ],
    "Use it right",
    [
        (0, "TEST the tester on a known source before AND after"),
        (0, "Match the detector's range to the system voltage"),
        (0, "Test-before-touch on every conductor, every time"),
        (0, "A 'no voltage' reading from an unproven tester proves NOTHING"),
    ],
    kicker="4.5 Testers",
    left_color=STEEL, right_color=AMBER_D,
)

bullets_slide(
    "Arc-Rated FR Clothing & Personal PPE",
    [
        (0, "Flame-resistant (FR) / arc-rated (AR) clothing protects against arc-flash thermal energy"),
        (1, "Rated in cal/cm² (ATPV / arc rating) — must meet or exceed the job's incident energy"),
        (1, "Self-extinguishing; will not melt and drip like synthetics (no untreated nylon/poly next to skin)"),
        (0, "Layering increases protection; outermost layer must be arc-rated"),
        (0, "Head/face/eyes: hard hat (Class E, electrical), arc-rated face shield/hood, safety glasses"),
        (0, "Hands: rubber gloves + leather protectors; feet: EH-rated footwear"),
        (0, "Hearing protection for arc-blast environments; no exposed conductive jewelry"),
        (0, "PPE is the LAST layer — it limits injury when other controls fail; it never replaces them"),
    ],
    kicker="4.6 FR / PPE",
)

table_slide(
    "PPE & Tool Inspection / Test Intervals (Typical)",
    ["Item", "Before each use", "Periodic electrical test"],
    [
        ["Rubber gloves", "Air & visual inspection", "Retest ≈ every 6 months"],
        ["Rubber sleeves", "Visual inspection", "Retest ≈ every 6 months"],
        ["Insulating blankets", "Visual inspection", "Retest ≈ every 12 months"],
        ["Line hose / covers", "Visual inspection", "When suspect / per program"],
        ["Hot sticks / live-line tools", "Wipe & visual inspection", "Periodic dielectric test (e.g., every 2 yrs)"],
        ["Protective grounds", "Visual; check clamps/strands", "Resistance/condition per program"],
        ["Voltage detectors", "Self-test before & after use", "Per manufacturer schedule"],
        ["FR clothing", "Inspect for damage/contamination", "Replace when worn/contaminated"],
    ],
    kicker="4.7 Intervals",
    col_widths=[1.7, 1.6, 1.9],
    fs=11.5,
    note="Intervals are TYPICAL and governed by OSHA 1910.137, ASTM F496/F478, and your employer's program. If equipment is dropped, suspect, or out of test date — remove it from service.",
)

add_module_figs(4)
bullets_slide(
    "Module 4 — Key Takeaways",
    [
        (0, "Rubber gloves (with leather protectors) are first-line insulation — selected by voltage class"),
        (0, "Air-test and inspect rubber goods before every use; honor retest dates"),
        (0, "Hot sticks isolate the worker; grounds protect against re-energization — both must be tested & sized"),
        (0, "Test your tester before AND after; test-before-touch every time"),
        (0, "Arc-rated FR PPE must meet or exceed the incident energy — and it's the last layer, not the plan"),
        (0, "Clean, dry, inspected, in-date — or out of service. No exceptions."),
    ],
    kicker="Module 4 Recap",
)

# ============================================================================
# MODULE 5 — FIELD APPLICATION
# ============================================================================
_state["module"] = "Module 5 — In the Field"
section_divider("5", "Putting It to Work in the Field",
                ["Pre-job planning",
                 "Hazard assessment / JSA",
                 "The tailboard / job briefing",
                 "Roles, responsibilities & communication",
                 "Emergency response & rescue"])

bullets_slide(
    "Pre-Job Planning — Before Boots Leave the Truck",
    [
        (0, "Understand the SCOPE: what work, on what equipment, at what voltage, energized or not"),
        (0, "Review prints, switching orders, clearances, and the circuit map"),
        (0, "Confirm the right crew, qualifications, and a competent person in charge"),
        (0, "Verify tools, PPE, cover-up, grounds — right ratings, inspected, in test date"),
        (0, "Check the weather: wind, lightning, heat, ice — conditions that stop or change the job"),
        (0, "Identify the access, the work zone, traffic control, and the public exposure"),
        (0, "Know the location: address, nearest cross-streets, gate codes — for emergency response"),
        (0, "Plan the rescue BEFORE the work — not after the incident"),
    ],
    kicker="5.1 Planning",
)

callout_slide(
    "Hazard Assessment / Job Safety Analysis (JSA)",
    [
        (0, "Break the job into steps; for each step identify the hazards and the controls"),
        (0, "Electrical hazards: energized parts, second points of contact, induced voltage, fault energy"),
        (0, "Mechanical/physical: falls, dropped objects, pinch points, equipment, rigging, tension"),
        (0, "Environmental: weather, traffic, terrain, wildlife, confined/excavation, public"),
        (0, "For each hazard, select the control from the hierarchy — eliminate where you can"),
        (0, "Reassess when ANYTHING changes: scope, weather, crew, equipment, the unexpected"),
    ],
    "info", "The four questions",
    ["What am I about to do?  What can hurt me?",
     "How do I control it?  What if the plan changes?",
     "A JSA is a living document — update it as the job evolves."],
    kicker="5.2 Hazard Assessment",
)

bullets_slide(
    "The Tailboard / Job Briefing — What & Why",
    [
        (0, "A tailboard is the crew safety briefing held before work starts (and again when things change)"),
        (0, "Required before each job; it aligns the whole crew on the plan and the hazards"),
        (0, "OSHA 1910.269(c) requires a job briefing covering at minimum:"),
        (1, "Hazards associated with the job"),
        (1, "Work procedures involved"),
        (1, "Special precautions"),
        (1, "Energy-source controls (clearances, grounds, LOTO)"),
        (1, "Personal protective equipment requirements"),
        (0, "More extensive briefing for complex/hazardous jobs; brief again after interruptions or changes"),
    ],
    kicker="5.3 Tailboard",
    lead="The tailboard is where the plan in your head becomes the plan in everyone's head.",
)

process_slide(
    "Running an Effective Tailboard",
    [
        ("1. Describe the job", "What are we doing, where, on what equipment, energized or de-energized, at what voltage."),
        ("2. Walk the hazards", "Go through the JSA: electrical (second points, MAD, induction), falls, traffic, weather, public."),
        ("3. State the controls", "Cover-up plan, grounds/EPZ, LOTO/clearance, MAD, PPE, traffic control — specific, not generic."),
        ("4. Assign roles", "Person in charge, who works, who tends, who grounds, who watches, who calls for help."),
        ("5. Confirm rescue & comms", "Location/address, nearest hospital, rescue plan & equipment, radio/phone, hand signals."),
        ("6. Engage the crew", "Ask questions, invite input, confirm understanding — everyone has stop-work authority."),
    ],
    kicker="5.3 Briefing",
)

two_col_slide(
    "Roles & Responsibilities on the Crew",
    "Key roles",
    [
        (0, "Person-in-Charge / Crew Leader — owns the plan & clearance"),
        (0, "Qualified worker(s) — perform the energized/skilled tasks"),
        (0, "Groundman / tender — supports, tends ropes, handles materials"),
        (0, "Safety watch / second person — eyes on the worker, ready to act"),
        (0, "Flagger / traffic control — protects the work zone & public"),
    ],
    "Shared responsibilities",
    [
        (0, "Everyone: know the plan and your part in it"),
        (0, "Everyone: maintain situational awareness"),
        (0, "Everyone has STOP-WORK authority — use it without fear"),
        (0, "Speak up: a question now beats an incident later"),
        (0, "No solo energized work — a qualified second person is present"),
    ],
    kicker="5.4 Roles",
    left_color=NAVY2, right_color=GREEN,
)

bullets_slide(
    "Communication — The Glue That Holds the Job Together",
    [
        (0, "Use clear, standardized language for switching and clearances — no ambiguity"),
        (0, "Three-way communication for critical commands: state → repeat back → confirm"),
        (0, "Confirm clearances and grounds with the system operator before and after"),
        (0, "Agree on hand signals for crane/bucket/rigging where radios are hard to hear"),
        (0, "Announce movements: 'coming down', 'tension on', 'clear' — keep the crew in the loop"),
        (0, "Re-brief after ANY interruption, crew change, or scope change"),
        (0, "If you're not sure what was said, STOP and confirm — assumptions kill"),
    ],
    kicker="5.5 Communication",
)

callout_slide(
    "Emergency Response & Rescue — Plan It Before You Need It",
    [
        (0, "Every energized/elevated job needs a rescue plan briefed in the tailboard"),
        (0, "Pole-top & bucket rescue: trained crew can lower an injured worker quickly"),
        (1, "Know the equipment (handline, rescue device) and practice the procedure"),
        (0, "Do NOT become the second victim — ensure the source of contact is removed/de-energized first"),
        (0, "Call for help early: know the address, cross-streets, and access"),
        (0, "CPR + AED for electrical contact / cardiac arrest — minutes matter; know where the AED is"),
        (0, "After contact, ALWAYS transport for medical evaluation — internal injury can be hidden"),
    ],
    "warn", "The first rule of rescue",
    ["Protect yourself first — a rescuer who gets electrocuted helps no one.",
     "De-energize / remove the source, then reach the patient.",
     "Speed saves lives, but only a SAFE rescue saves two."],
    kicker="5.6 Rescue",
)

bullets_slide(
    "Weather, Environment & Public Safety",
    [
        (0, "Lightning / approaching storms: stop work and clear elevated/exposed positions"),
        (0, "High wind: affects conductors, booms, rigging, and footing — know your limits"),
        (0, "Wet conditions: drastically reduce body & equipment insulation — reassess hot work"),
        (0, "Heat/cold: fatigue and impaired judgment are hazards in themselves — hydrate, rotate"),
        (0, "Ice/snow loading: added conductor weight and slip hazards"),
        (0, "Public protection: barricades, signage, traffic control; keep bystanders clear of downed lines"),
        (0, "Wildlife & vegetation: animals cause faults; trees create clearance & contact hazards"),
    ],
    kicker="5.7 Environment",
)

add_module_figs(5)
bullets_slide(
    "Module 5 — Key Takeaways",
    [
        (0, "Plan the job — and the rescue — before boots leave the truck"),
        (0, "A JSA pairs every step with its hazards and controls, and updates when things change"),
        (0, "The tailboard puts the plan in everyone's head: hazards, procedures, controls, PPE, roles"),
        (0, "Define roles and use disciplined, three-way communication"),
        (0, "Everyone has stop-work authority — and the obligation to use it"),
        (0, "Rescue first protects the rescuer; after any contact, get medical evaluation"),
    ],
    kicker="Module 5 Recap",
)

# ============================================================================
# MODULE 6 — REGULATIONS & STANDARDS
# ============================================================================
_state["module"] = "Module 6 — Regulations & Standards"
section_divider("6", "Regulations & Standards",
                ["OSHA 1910.269 & 1926 Subpart V",
                 "The National Electrical Safety Code (NESC)",
                 "ASTM & IEEE equipment standards",
                 "Qualification & apprenticeship",
                 "Recordkeeping & program duties"])

bullets_slide(
    "Why Regulations Exist — and Why You Should Know Them",
    [
        (0, "Every safe-work rule in this course is backed by a regulation or consensus standard"),
        (0, "Regulations set the MINIMUM — your employer's program may be stricter (follow the stricter)"),
        (0, "Knowing the 'why' and the 'where it's written' makes you a better, safer, more credible worker"),
        (0, "The big three for line work:"),
        (1, "OSHA — federal law: how the work must be done safely (enforceable)"),
        (1, "NESC (IEEE C2) — how the system is built, spaced, and maintained"),
        (1, "ASTM / IEEE / ANSI — how the tools and PPE are made, rated, and tested"),
        (0, "Compliance is the floor; the goal is everyone home safe — which often means going beyond the floor"),
    ],
    kicker="6.1 Why Regs",
)

bullets_slide(
    "OSHA 1910.269 — Electric Power Generation, Transmission & Distribution",
    [
        (0, "The primary OSHA standard for operating & maintaining electric power systems (general industry)"),
        (0, "Key provisions every lineworker should know:"),
        (1, "(a) Training & qualified-employee requirements"),
        (1, "(c) Job briefings (the tailboard) — required content and frequency"),
        (1, "(l) Working on or near energized parts; minimum approach distances (Tables R-3 / R-6)"),
        (1, "(m) De-energizing lines & equipment"),
        (1, "(n) Grounding for the protection of employees"),
        (1, "Arc-flash hazard assessment and arc-rated apparel requirements"),
        (0, "Companion construction standard: 29 CFR 1926 Subpart V (power transmission & distribution)"),
    ],
    kicker="6.2 OSHA 1910.269",
    lead="If you read one regulation cover to cover, make it this one.",
)

two_col_slide(
    "OSHA 1910.269 vs 1926 Subpart V",
    "1910.269 (General Industry)",
    [
        (0, "Operation & MAINTENANCE of existing power systems"),
        (0, "Day-to-day line & substation work"),
        (0, "Job briefings, MAD, de-energizing, grounding"),
        (0, "Harmonized with Subpart V on key requirements"),
    ],
    "1926 Subpart V (Construction)",
    [
        (0, "CONSTRUCTION of power transmission & distribution"),
        (0, "Building new lines and major rebuilds"),
        (0, "Mirrors 1910.269 on MAD, briefings, grounding"),
        (0, "Use the standard that matches the type of work"),
    ],
    kicker="6.2 Scope",
    left_color=NAVY2, right_color=STEEL,
)

bullets_slide(
    "The National Electrical Safety Code (NESC / IEEE C2)",
    [
        (0, "NESC governs the safe installation, operation & maintenance of utility electric & communication systems"),
        (0, "Distinct from the NEC (NFPA 70), which covers premises/building wiring"),
        (0, "Sets rules for the SYSTEM that protect both workers and the public, including:"),
        (1, "Conductor clearances — ground, structures, other lines, and to the public"),
        (1, "Strength & loading of structures; grades of construction"),
        (1, "Grounding methods and worker safety rules (Part 4 covers work rules)"),
        (0, "Updated on a multi-year cycle; the edition in force depends on jurisdiction"),
        (0, "When clearances or construction look wrong in the field — NESC is the reference"),
    ],
    kicker="6.3 NESC",
)

table_slide(
    "Key ASTM / IEEE / ANSI Equipment Standards",
    ["Standard", "Covers"],
    [
        ["ASTM D120", "Rubber insulating GLOVES — manufacture & rating"],
        ["ASTM D1051", "Rubber insulating SLEEVES"],
        ["ASTM D1048", "Rubber insulating BLANKETS"],
        ["ASTM D1049 / D1050", "Insulating COVERS / LINE HOSE"],
        ["ASTM F496", "In-service CARE of rubber gloves & sleeves (test intervals)"],
        ["ASTM F478", "In-service care of insulating LINE HOSE & covers"],
        ["ASTM F711", "Fiberglass rod & tube for LIVE-LINE TOOLS"],
        ["ASTM F855", "Temporary protective GROUNDS (clamps & cables)"],
        ["ASTM F1506 / NFPA 70E", "Arc-rated FR CLOTHING & arc-flash protection"],
        ["IEEE C2 (NESC)", "Utility system safety code"],
        ["ANSI Z89.1", "Hard hats (electrical Class E)"],
    ],
    kicker="6.4 ASTM/IEEE",
    col_widths=[1.4, 3.4],
    fs=11.5,
    note="These standards define how your PPE and tools are built, rated, and re-tested — the numbers on your gloves and grounds trace back here.",
)

bullets_slide(
    "Qualified Workers, Training & Apprenticeship",
    [
        (0, "OSHA requires that only QUALIFIED employees work on or near energized parts"),
        (0, "Qualified = trained AND able to demonstrate skills/knowledge for the construction & operation involved"),
        (1, "Knows how to distinguish energized parts, determine voltages, and maintain MAD"),
        (0, "Lineworker apprenticeships (often ~3–4 years, thousands of OJT hours) build that qualification"),
        (1, "Classroom theory + supervised on-the-job training + skills verification"),
        (0, "Ongoing training: refreshers, new procedures, CPR/AED & pole-top rescue, equipment-specific"),
        (0, "Document training and re-qualification — it is both a legal requirement and proof of competence"),
    ],
    kicker="6.5 Qualification",
)

two_col_slide(
    "Recordkeeping & Program Responsibilities",
    "Records to keep",
    [
        (0, "Training & qualification records"),
        (0, "PPE/tool inspection & electrical test dates"),
        (0, "Job briefings / tailboards (per program)"),
        (0, "Switching orders, clearances & grounding records"),
        (0, "Incident & near-miss reports and investigations"),
    ],
    "Employer program duties",
    [
        (0, "Written safe-work & electrical safety program"),
        (0, "Hazard/arc-flash assessments & MAD tables"),
        (0, "Provide & maintain rated PPE and tools"),
        (0, "Enforce procedures; support stop-work authority"),
        (0, "Investigate, learn, and improve after every event"),
    ],
    kicker="6.6 Records",
    left_color=NAVY2, right_color=GREEN,
)

bullets_slide(
    "Module 6 — Key Takeaways",
    [
        (0, "Regulations set the minimum; your program may be stricter — always follow the stricter rule"),
        (0, "OSHA 1910.269 (and 1926 Subpart V) govern how the work is done — briefings, MAD, LOTO, grounding"),
        (0, "NESC (IEEE C2) governs how the system is built, spaced, and maintained"),
        (0, "ASTM/IEEE/ANSI define how your PPE & tools are rated and tested"),
        (0, "Only qualified workers do energized work — earned through apprenticeship & ongoing training"),
        (0, "Keep records: training, test dates, briefings, clearances, incidents"),
    ],
    kicker="Module 6 Recap",
)

# ============================================================================
# MODULE 7 — ARC-FLASH, FIRST AID & RESCUE
# ============================================================================
_state["module"] = "Module 7 — Arc-Flash, First Aid & Rescue"
section_divider("7", "Arc-Flash Analysis, First Aid & Rescue",
                ["What drives incident energy",
                 "The arc-flash boundary & working distance",
                 "PPE categories & the label",
                 "Electrical injury & burns",
                 "Scene control, CPR & AED"])

bullets_slide(
    "Arc-Flash Analysis — What Drives Incident Energy",
    [
        (0, "Incident energy is the thermal energy reaching a surface a set distance from an arc — measured in cal/cm²"),
        (0, "It rises with the available fault (bolted) current and with how long the arc lasts"),
        (1, "Clearing time is decisive: a faster-tripping device means far less energy"),
        (0, "It falls with the SQUARE of the distance from the arc"),
        (0, "An analysis (IEEE 1584 / NFPA 70E) calculates it for each location"),
        (1, "Inputs: fault current, electrode gap & configuration, equipment type, working distance"),
        (0, "The result drives the required PPE arc rating and the arc-flash boundary"),
    ],
    kicker="7.1 Incident Energy",
    lead="Why two jobs at the same voltage can carry wildly different arc-flash hazard.",
)

callout_slide(
    "The Arc-Flash Boundary & Working Distance",
    [
        (0, "Working distance: the arc-to-worker distance used in the calc (often 18 in for distribution)"),
        (1, "The incident energy at that distance sets the PPE arc rating you must wear"),
        (0, "Arc-flash boundary: the distance where incident energy = 1.2 cal/cm²"),
        (1, "1.2 cal/cm² is the onset of a second-degree burn to bare skin"),
        (0, "Anyone inside the boundary must wear arc-rated PPE"),
        (0, "Shock approach boundaries (limited / restricted) are separate and based on voltage"),
    ],
    "info", "Two different boundaries",
    ["Arc-flash boundary = thermal (cal/cm²).",
     "Shock approach boundaries = voltage-based (the MAD family).",
     "Respect both — they guard against different hazards."],
    kicker="7.2 Boundary & Distance",
)

table_slide(
    "Arc-Flash PPE Categories (NFPA 70E)",
    ["Category", "Min arc rating", "Typical ensemble"],
    [
        ["CAT 1", "≥ 4 cal/cm²", "AR shirt & pants, hard hat, arc face shield, gloves"],
        ["CAT 2", "≥ 8 cal/cm²", "Above + arc-rated hood or balaclava with face shield"],
        ["CAT 3", "≥ 25 cal/cm²", "AR flash-suit jacket & trousers + arc-rated hood"],
        ["CAT 4", "≥ 40 cal/cm²", "Heavy AR flash suit & hood — full ensemble"],
    ],
    kicker="7.3 PPE Categories",
    col_widths=[1.0, 1.4, 3.2],
    note="The category method is one approach; an incident-energy analysis is the other. Either way the garment's arc rating (ATPV) must meet or exceed the incident energy. Above ~40 cal/cm² the blast hazard alone may make energized work unacceptable — de-energize.",
)

bullets_slide(
    "Reading the Arc-Flash Label",
    [
        (0, "Equipment should carry a label produced from the facility's arc-flash analysis"),
        (0, "Typical fields:"),
        (1, "Nominal system voltage · arc-flash boundary · incident energy @ working distance"),
        (1, "Required PPE category · minimum glove class · shock approach boundaries"),
        (0, "Use it to select PPE and set your boundaries BEFORE energizing or opening equipment"),
        (0, "Missing, outdated, or illegible label → STOP and get the analysis"),
        (0, "Labels change when the system changes — new transformer, settings, or fault current"),
    ],
    kicker="7.4 The Label",
    lead="The label turns the engineering study into a field decision.",
)

bullets_slide(
    "Electrical Injury — What Happens to the Body",
    [
        (0, "Cardiac: ventricular fibrillation or arrest from current across the chest"),
        (0, "Burns: surface arc burns AND internal burns along the current path"),
        (1, "Small entry/exit wounds can hide massive internal tissue damage"),
        (0, "Respiratory: chest muscles lock — breathing can stop"),
        (0, "Neurological: nerve damage, confusion, lasting deficits"),
        (0, "Trauma: falls from the shock, blast-thrown objects and molten metal"),
        (0, "Delayed effects: kidney injury from tissue breakdown; rhythm problems hours later"),
    ],
    kicker="7.5 Electrical Injury",
    lead="One contact can injure several body systems at once — which is why every contact gets evaluated.",
)

callout_slide(
    "Emergency Response — Scene Control & the Rescuer",
    [
        (0, "Protect yourself FIRST — a second victim helps no one"),
        (0, "Do NOT touch the patient until the source is removed / de-energized and you know it"),
        (0, "Treat downed conductors and the ground around them as energized — keep a wide barrier"),
        (0, "Call for help early with a precise location; send someone for the AED"),
        (0, "Decide who is in charge of the rescue and communicate clearly"),
    ],
    "warn", "The order is fixed",
    ["1. Make the scene safe (remove the source).",
     "2. Reach the patient.",
     "3. Treat — CPR/AED, then burns.",
     "Never skip step 1 to save time."],
    kicker="7.6 Scene Control",
)

process_slide(
    "CPR & AED for Electrical Contact",
    [
        ("1. Scene safe", "Confirm the source is removed / de-energized before you approach the patient."),
        ("2. Check & call", "Check responsiveness and breathing; call 911 and send for the AED."),
        ("3. Compressions", "Center of the chest, 100–120 per minute, about 2 inches deep, full recoil."),
        ("4. Apply the AED", "Power on, attach pads, follow the voice prompts; clear for analysis and shock."),
        ("5. Continue", "Resume compressions immediately; cycle until EMS takes over."),
    ],
    kicker="7.7 CPR & AED",
    lead="Push hard, push fast, minimize interruptions.",
)

bullets_slide(
    "Burns & Ongoing Care",
    [
        (0, "Stop the burning; cool a thermal burn with clean water where appropriate"),
        (0, "Cover with a clean, dry dressing — do not break blisters or apply ointments to severe burns"),
        (0, "Treat for shock; keep the patient warm"),
        (0, "Electrical burns are deceptive — internal damage far exceeds the visible wound"),
        (0, "ALWAYS transport for medical evaluation after any contact, even if they feel fine"),
        (0, "Report and investigate the contact — every event is a chance to prevent the next"),
    ],
    kicker="7.8 Burns & Care",
)

add_module_figs(7)
bullets_slide(
    "Module 7 — Key Takeaways",
    [
        (0, "Incident energy (cal/cm²) rises with fault current & clearing time, falls with distance²"),
        (0, "The arc-flash boundary is where energy = 1.2 cal/cm²; PPE arc rating must meet or exceed the incident energy"),
        (0, "Read the label and set your boundaries before opening equipment"),
        (0, "One contact injures many systems — cardiac, burns, neuro, trauma; effects can be delayed"),
        (0, "Rescue order is fixed: make the scene safe, reach the patient, then treat"),
        (0, "CPR + AED for arrest; transport every contact for medical evaluation"),
    ],
    kicker="Module 7 Recap",
)

# ============================================================================
# CLOSE
# ============================================================================
_state["module"] = "Course Wrap-Up"
section_divider("✓", "Bringing It All Together",
                ["The mental model from theory to the field",
                 "Daily safe-work disciplines",
                 "Final reminders",
                 "Resources & references"])

process_slide(
    "From Theory to the Field — The Mental Model",
    [
        ("Understand the energy", "Voltage pushes current; current through you kills. Know your voltage class and fault energy."),
        ("Find the circuit", "Where would current flow through ME? Locate the SECOND POINT OF CONTACT — every time."),
        ("Choose the control", "De-energize if you can. If hot: isolate AND insulate — MAD, cover-up, gloves, grounds, EPZ."),
        ("Verify", "Test the tester, test-before-touch, prove dead, apply grounds, inspect every piece of PPE."),
        ("Brief & assign", "Tailboard the hazards, controls, roles, comms, and rescue. Everyone understands. Everyone can stop work."),
        ("Work & reassess", "Maintain controls all the way through. When ANYTHING changes — stop, re-brief, re-assess."),
    ],
    kicker="Synthesis",
)

bullets_slide(
    "Ten Disciplines That Keep Crews Alive",
    [
        (0, "Treat every conductor as energized until YOU prove it dead and grounded"),
        (0, "Always identify and eliminate the second point of contact"),
        (0, "Maintain minimum approach distance — for your body, tools, and equipment"),
        (0, "Cover up: nearest hazard first, remove in reverse"),
        (0, "Test the tester before and after; test-before-touch"),
        (0, "Ground and bond an equipotential zone at the work site"),
        (0, "Inspect rubber goods and tools before every use; honor test dates"),
        (0, "Hold a real tailboard — and re-brief when anything changes"),
        (0, "Plan the rescue before the work; never become the second victim"),
        (0, "Use your stop-work authority — no job is worth your life"),
    ],
    kicker="Daily Discipline",
    lead="If you remember nothing else, remember these.",
)

callout_slide(
    "Final Word",
    [
        (0, "Electricity is invisible, silent, and unforgiving — but it is also entirely predictable"),
        (0, "It follows the same laws every time, which means the hazard is always controllable"),
        (0, "The theory in Module 1 is not academic — it is WHY every rule in Modules 3–6 works"),
        (0, "Skills fade and complacency creeps in — that is why we train, brief, and reassess relentlessly"),
        (0, "Look out for the worker next to you as carefully as you look out for yourself"),
    ],
    "safe", "The whole point",
    ["Everyone goes home, every day.",
     "That is the only acceptable outcome of any job — and it is achievable on every job."],
    kicker="Closing",
)

two_col_slide(
    "Resources & References",
    "Regulations & codes",
    [
        (0, "OSHA 29 CFR 1910.269 — GT&D (general industry)"),
        (0, "OSHA 29 CFR 1926 Subpart V — construction"),
        (0, "OSHA 1910.137 — electrical protective equipment"),
        (0, "OSHA 1910.147 — control of hazardous energy (LOTO)"),
        (0, "NESC — IEEE C2; NFPA 70E — arc-flash safety"),
    ],
    "Standards & training",
    [
        (0, "ASTM F18 committee standards (D120, F496, F855, etc.)"),
        (0, "IEEE standards for live-line work & grounding"),
        (0, "Your employer's safe-work practices & MAD tables"),
        (0, "Apprenticeship program materials & manufacturer manuals"),
        (0, "First aid / CPR / AED & pole-top rescue certification"),
    ],
    kicker="References",
    left_color=NAVY2, right_color=STEEL,
)

# Final thank-you / Q&A
s = slide()
grad_rect(s, 0, 0, 13.333, 7.5, NAVY, (0x06, 0x16, 0x2B), angle=60)
rect(s, 0.0, 0.0, 13.333, 0.14, fill=AMBER)
rect(s, 0.9, 2.5, 0.18, 2.0, fill=AMBER)
tb, tf = textbox(s, 1.3, 2.55, 11.0, 2.5)
para(tf, "Questions & Discussion", size=46, color=WHITE, bold=True, first=True, after=10)
para(tf, "Review the hazards on YOUR system, with YOUR equipment, against YOUR procedures.",
     size=18, color=RGBColor(0xC8, 0xD6, 0xE6), line=1.2)
para(tf, "Stay sharp. Watch your second point of contact. Everyone goes home.",
     size=16, color=AMBER, italic=True, before=14)
rect(s, 0.0, 7.2, 13.333, 0.06, fill=AMBER)
_state["n"] += 1

# ----------------------------------------------------------------------------
# Attach speaker notes (match each slide's title text to notes.json)
notes_applied = 0
for _sld in prs.slides:
    if _sld.has_notes_slide and _sld.notes_slide.notes_text_frame.text.strip():
        continue  # image slides already carry their caption as notes
    _done = False
    for _shp in _sld.shapes:
        if _done or not _shp.has_text_frame:
            continue
        _cands = [_shp.text_frame.text]
        if _shp.text_frame.paragraphs:
            _cands.append(_shp.text_frame.paragraphs[0].text)
        for _c in _cands:
            _key = _norm(_c)
            if _key in NOTES:
                _sld.notes_slide.notes_text_frame.text = NOTES[_key]
                notes_applied += 1
                _done = True
                break

out = "Powerline_Electrical_Theory_and_Safety_Training.pptx"
prs.save(out)
print("Saved", out, "with", len(prs.slides._sldIdLst), "slides;",
      notes_applied, "title-matched notes +", sum(len(v) for v in MODULE_FIGS.values()), "diagram-caption notes")
